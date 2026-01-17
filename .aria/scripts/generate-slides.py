#!/usr/bin/env python3
"""
ARIA Slide Generation Script

Two paths:
1. NotebookLM (requires notebooklm-py + Google auth)
2. Local pptx (requires python-pptx)

Usage:
    python generate-slides.py --focus .aria/outputs/FOCUS.md \
                              --idea .aria/docs/IDEA.md \
                              --sources paper.pdf \
                              --method nblm|pptx
"""

import argparse
import asyncio
import os
import re
from datetime import datetime
from pathlib import Path

# Constants
ARIA_DIR = Path('.aria')
OUTPUTS_DIR = ARIA_DIR / 'outputs'

FOCUS_PROMPT = """Analyze all sources and provide a structured synthesis:
1. The Core: Identify the top 3 foundational elements (concepts, entities, or arguments) that are absolutely required to understand this corpus. Define them briefly.
2. The Synthesis: Coalesce the findings by listing the top 5-10 unifying ideas or themes that connect the "Core Trinity" together. Explain how these ideas turn the separate elements into a cohesive whole."""

SLIDES_PROMPT = """Intent: explain in detail the key ideas from these docs - USE THE [FOCUS doc] as the guide to highlight each important aspect we need to bring forth and explain

Must DO: provide detailed learning deck to explain the workflow - be verbose - use unconventional spatial and verbal slide design techniques rooted in cognitive science for maximum learning - make this a long deck 20 plus slides if necessary. Use charts, graphs, flow diagrams, and other visuals liberally to get the message across. Break concepts down for ease of intake for new learners or people unfamiliar with the content - again use diagrams liberally - capture main steps in all workflows - ensure the high level process is clear. Provide a clear concise view at the end of complete flow"""


def parse_focus_doc(focus_path: Path) -> dict:
    """Parse FOCUS.md to extract Core and Synthesis sections."""
    content = focus_path.read_text()

    result = {
        'core_ideas': [],
        'synthesis_matrix': [],
        'raw': content
    }

    # Extract The Core section (matches "The Core", "Core Trinity", "Core Ideas", etc.)
    core_match = re.search(
        r'(?:The Core|Core Trinity|Core Ideas?|Foundational Elements?).*?(?=The Synthesis|Synthesis|$)',
        content,
        re.IGNORECASE | re.DOTALL
    )
    if core_match:
        ideas = re.findall(r'^\s*[-*\d.]+\s*(.+)$', core_match.group(), re.MULTILINE)
        result['core_ideas'] = [i.strip() for i in ideas[:5]]

    # Extract The Synthesis section (matches "The Synthesis", "Synthesis Matrix", etc.)
    synthesis_match = re.search(
        r'(?:The Synthesis|Synthesis Matrix?|Unifying (?:Ideas|Themes)).*',
        content,
        re.IGNORECASE | re.DOTALL
    )
    if synthesis_match:
        themes = re.findall(r'^\s*[-*\d.]+\s*(.+)$', synthesis_match.group(), re.MULTILINE)
        result['synthesis_matrix'] = [t.strip() for t in themes[:10]]

    return result


def extract_title_from_idea(idea_path: Path) -> str:
    """Extract title from IDEA.md."""
    content = idea_path.read_text()
    
    # Try to find # Title
    title_match = re.search(r'^#\s+(.+)$', content, re.MULTILINE)
    if title_match:
        return title_match.group(1).strip()
    
    # Fallback to first line
    first_line = content.strip().split('\n')[0]
    return first_line[:50].strip('# ')


# ============================================================
# NotebookLM Path
# ============================================================

async def generate_nblm(focus_path: Path, idea_path: Path, source_paths: list) -> Path:
    """Generate slides via NotebookLM."""
    try:
        from notebooklm import NotebookLMClient
    except ImportError:
        print("ERROR: notebooklm-py not installed")
        print("Install with: pip install \"notebooklm-py[browser]\"")
        print("Then run: playwright install chromium")
        print("Then run: notebooklm login")
        print("Falling back to pptx...")
        return await generate_pptx(focus_path, idea_path, source_paths)

    title = extract_title_from_idea(idea_path)
    timestamp = datetime.now().strftime('%Y%m%d-%H%M')

    try:
        async with await NotebookLMClient.from_storage() as client:
            # ALWAYS create a new notebook - never search for existing
            # Each research session should have its own notebook for traceability

            # Create notebook
            print(f"Creating NotebookLM notebook: {title}")
            notebook = await client.notebooks.create(f"Slides: {title}")
            print(f"  Notebook ID: {notebook.id}")

            # Add sources - use add_url for text content or add_file for files
            print("Adding FOCUS.md content...")
            # Write temp file for FOCUS content
            focus_temp = OUTPUTS_DIR / f"_temp_focus_{timestamp}.md"
            focus_temp.write_text(focus_path.read_text())
            await client.sources.add_file(notebook.id, str(focus_temp))

            print("Adding IDEA.md content...")
            idea_temp = OUTPUTS_DIR / f"_temp_idea_{timestamp}.md"
            idea_temp.write_text(idea_path.read_text())
            await client.sources.add_file(notebook.id, str(idea_temp))

            for source in source_paths:
                source_path = Path(source)
                if source_path.exists():
                    print(f"Adding {source_path.name}...")
                    await client.sources.add_file(notebook.id, str(source_path))

            # Send slide generation prompt
            print("\nSending slide generation prompt...")
            print(f"  Prompt: {SLIDES_PROMPT[:100]}...")
            response = await client.chat.ask(notebook.id, SLIDES_PROMPT)
            print(f"  Prompt sent successfully")

            # Start slide deck generation (don't wait - takes 5-10 min)
            print("\nStarting slide deck generation...")
            print("  NOTE: This takes 5-10 minutes. Check NotebookLM directly.")
            status = await client.artifacts.generate_slide_deck(notebook.id)
            print(f"  Task started: {status.task_id}")
            print(f"  Status: {status.status}")

            # Get notebook URL
            notebook_url = f"https://notebooklm.google.com/notebook/{notebook.id}"

            # Cleanup temp files
            focus_temp.unlink(missing_ok=True)
            idea_temp.unlink(missing_ok=True)

            print("\n" + "="*60)
            print("SLIDE GENERATION STARTED")
            print("="*60)
            print(f"\nNotebook URL: {notebook_url}")
            print("\nNext steps:")
            print("  1. Open the URL above in your browser")
            print("  2. Wait 5-10 minutes for slide deck to generate")
            print("  3. Download from NotebookLM when ready")
            print("="*60)

            return notebook_url

    except Exception as e:
        error_str = str(e).lower()
        print("\n" + "="*60)
        print("NOTEBOOKLM FAILURE - DIAGNOSTIC REPORT")
        print("="*60)
        print(f"\nError: {e}")

        # Diagnose specific failure type
        if 'auth' in error_str or 'login' in error_str or 'credential' in error_str or 'token' in error_str:
            print("\n" + "-"*40)
            print("DIAGNOSIS: LOGIN/AUTHENTICATION FAILURE")
            print("-"*40)
            print("\nCause: Google authentication expired or invalid")
            print("\nFix:")
            print("  1. Run: notebooklm login")
            print("  2. Complete Google sign-in in browser")
            print("  3. Re-run slide generation")

        elif 'upload' in error_str or 'file' in error_str or 'source' in error_str or 'add' in error_str:
            print("\n" + "-"*40)
            print("DIAGNOSIS: FILE UPLOAD FAILURE")
            print("-"*40)
            print("\nCause: Failed to upload sources to NotebookLM")
            print("\nCheck:")
            print(f"  - FOCUS.md exists: {focus_path.exists()}")
            print(f"  - IDEA.md exists: {idea_path.exists()}")
            print("  - Files contain ASCII-only characters (no Unicode)")
            print("  - File sizes are under NotebookLM limits")
            print("\nFix:")
            print("  1. Verify files exist and are readable")
            print("  2. Check for Unicode characters (box-drawing, arrows)")
            print("  3. Regenerate FOCUS.md with ASCII-only format")

        elif 'slide' in error_str or 'deck' in error_str or 'generate' in error_str or 'artifact' in error_str:
            print("\n" + "-"*40)
            print("DIAGNOSIS: SLIDE CREATION FAILURE")
            print("-"*40)
            print("\nCause: NotebookLM failed to start slide generation")
            print("\nPossible reasons:")
            print("  - NotebookLM service temporarily unavailable")
            print("  - Sources didn't process correctly")
            print("  - Rate limiting")
            print("\nFix:")
            print("  1. Check notebook directly: https://notebooklm.google.com")
            print("  2. Verify sources appear in the notebook")
            print("  3. Try manual slide generation from NotebookLM UI")

        elif 'notebook' in error_str or 'create' in error_str:
            print("\n" + "-"*40)
            print("DIAGNOSIS: NOTEBOOK CREATION FAILURE")
            print("-"*40)
            print("\nCause: Could not create new notebook")
            print("\nFix:")
            print("  1. Check NotebookLM is accessible: https://notebooklm.google.com")
            print("  2. Verify Google account has NotebookLM access")
            print("  3. Run: notebooklm login")

        else:
            print("\n" + "-"*40)
            print("DIAGNOSIS: UNKNOWN ERROR")
            print("-"*40)
            print("\nCould not determine specific failure type.")

        print("\n" + "="*60)
        print("RECOVERY OPTIONS")
        print("="*60)
        print("\n  [1] Fix issue and re-run: python .aria/scripts/generate-slides.py --method nblm")
        print("  [2] Check NotebookLM directly: https://notebooklm.google.com")
        print("  [3] Re-authenticate: notebooklm login")
        print("  [4] Use fallback pptx: python .aria/scripts/generate-slides.py --method pptx")
        print("\nDo NOT auto-retry. User must diagnose and decide next step.")
        print("="*60)
        return None


# ============================================================
# pptx Fallback Path
# ============================================================

async def generate_pptx(focus_path: Path, idea_path: Path, source_paths: list) -> Path:
    """Generate slides via python-pptx."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RgbColor
    except ImportError:
        print("ERROR: python-pptx not installed")
        print("Install with: pip install python-pptx")
        raise SystemExit(1)
    
    title = extract_title_from_idea(idea_path)
    focus_data = parse_focus_doc(focus_path)
    idea_content = idea_path.read_text()
    timestamp = datetime.now().strftime('%Y%m%d-%H%M')
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)
    
    # Helper functions
    def add_title_slide(title_text, subtitle_text=""):
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(44)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        if subtitle_text:
            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12), Inches(1))
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle_text
            p.font.size = Pt(24)
            p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_content_slide(title_text, bullets, subtitle=""):
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(32)
        p.font.bold = True
        
        # Subtitle
        y_offset = 1.1
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_offset), Inches(12), Inches(0.5))
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(18)
            p.font.italic = True
            y_offset = 1.6
        
        # Bullets
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_offset), Inches(12), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(20)
            p.space_after = Pt(12)
        
        return slide
    
    def add_diagram_placeholder(title_text, description):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(32)
        p.font.bold = True
        
        # Placeholder box
        shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(1), Inches(1.5), Inches(11), Inches(5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RgbColor(240, 240, 240)
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(10), Inches(1))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"[DIAGRAM: {description}]"
        p.font.size = Pt(18)
        p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    # Build presentation
    print("Building pptx presentation...")
    
    # 1. Title slide
    add_title_slide(title, "Research Synthesis Deck")
    
    # 2. Overview slide
    add_content_slide(
        "Overview",
        ["This deck synthesizes key research findings",
         f"Based on {len(focus_data['core_ideas'])} core ideas",
         f"Connected by {len(focus_data['synthesis_matrix'])} unifying themes"],
        "What you'll learn"
    )
    
    # 3. Core Ideas overview
    add_content_slide(
        "Core Ideas",
        focus_data['core_ideas'] or ["Core ideas extracted from Focus document"],
        "Foundational elements required to understand this topic"
    )
    
    # 4. Individual Core Idea slides
    for i, idea in enumerate(focus_data['core_ideas'][:3], 1):
        add_content_slide(
            f"Core Idea {i}",
            [idea, "Key implications:", "• [Detail from source documents]"],
            "Deep dive"
        )
        add_diagram_placeholder(
            f"Core Idea {i}: Visual",
            f"Diagram illustrating: {idea[:50]}..."
        )
    
    # 5. Synthesis Matrix
    add_content_slide(
        "Synthesis Matrix",
        focus_data['synthesis_matrix'][:5] or ["Themes connecting core ideas"],
        "Unifying ideas that create a cohesive whole"
    )
    
    # 6. Individual theme slides
    for i, theme in enumerate(focus_data['synthesis_matrix'][:5], 1):
        add_content_slide(
            f"Theme {i}",
            [theme, "How it connects:", "• [Connection details]"],
            "Synthesis"
        )
    
    # 7. Workflow overview
    add_diagram_placeholder(
        "Complete Workflow",
        "High-level process flow from inputs to outputs"
    )
    
    # 8. Summary slide
    summary_points = [
        f"Core Ideas: {', '.join(focus_data['core_ideas'][:3]) or 'See Focus doc'}",
        f"Key Themes: {len(focus_data['synthesis_matrix'])} unifying concepts",
        "Next Steps: Review source materials for depth"
    ]
    add_content_slide("Summary", summary_points, "Key takeaways")
    
    # 9. Final slide
    add_title_slide("Questions?", f"Generated {datetime.now().strftime('%Y-%m-%d')}")
    
    # Save
    output_path = OUTPUTS_DIR / f"slides-{timestamp}.pptx"
    OUTPUTS_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    
    print(f"Slides saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
    print("\nNOTE: pptx contains placeholder diagrams.")
    print("For richer visuals, use NotebookLM path.")
    
    return output_path


# ============================================================
# Main
# ============================================================

async def main():
    parser = argparse.ArgumentParser(description='Generate slides from research artifacts')
    parser.add_argument('--focus', type=Path, default=OUTPUTS_DIR / 'FOCUS.md',
                        help='Path to FOCUS.md')
    parser.add_argument('--idea', type=Path, default=ARIA_DIR / 'docs' / 'IDEA.md',
                        help='Path to IDEA.md')
    parser.add_argument('--sources', nargs='*', default=[],
                        help='Additional source files (PDFs, MDs)')
    parser.add_argument('--method', choices=['nblm', 'pptx'], default='pptx',
                        help='Generation method: nblm (NotebookLM) or pptx')
    
    args = parser.parse_args()
    
    # Validate inputs
    if not args.focus.exists():
        print(f"ERROR: Focus doc not found: {args.focus}")
        print("Run Focus generation first, or create FOCUS.md manually.")
        raise SystemExit(1)
    
    if not args.idea.exists():
        print(f"ERROR: IDEA.md not found: {args.idea}")
        raise SystemExit(1)
    
    # Generate
    if args.method == 'nblm':
        output = await generate_nblm(args.focus, args.idea, args.sources)
    else:
        output = await generate_pptx(args.focus, args.idea, args.sources)
    
    print(f"\nDone! Output: {output}")


if __name__ == '__main__':
    asyncio.run(main())
